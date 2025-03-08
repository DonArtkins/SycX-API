import os
from transformers import pipeline
import torch
from PIL import Image
import pytesseract
from pdfminer.high_level import extract_text as extract_pdf_text
from docx import Document 
from openpyxl import load_workbook
from pptx import Presentation
import io
import nltk
from flask import current_app
import logging

class FileProcessor:
    def __init__(self):
        # Initialize the summarization model with configuration from app config
        os.environ['MALLOC_TRIM_THRESHOLD_'] = str(current_app.config['MALLOC_TRIM_THRESHOLD'])
        os.environ['PYTORCH_CUDA_ALLOC_CONF'] = current_app.config['PYTORCH_CUDA_ALLOC_CONF']
        
        # Direct file upload to HuggingFace
        self.summarizer = pipeline(
            "summarization",
            model="facebook/bart-large-cnn",
            device=-1 # Use CPU to ensure stability
        )

        nltk.download('punkt', quiet=True)

    def _optimize_length_params(self, text_length, summary_depth):
        """Calculate optimal length parameters based on summary depth"""
        # Add more granular depth configs
        depth_configs = {
            0.0: {'max_ratio': 0.05, 'min_ratio': 0.02, 'title_length': 2},  # Very concise
            1.0: {'max_ratio': 0.15, 'min_ratio': 0.05, 'title_length': 3},  # Concise
            2.0: {'max_ratio': 0.30, 'min_ratio': 0.10, 'title_length': 4},  # Balanced
            3.0: {'max_ratio': 0.40, 'min_ratio': 0.20, 'title_length': 5},  # Detailed  
            4.0: {'max_ratio': 0.60, 'min_ratio': 0.30, 'title_length': 6}   # Very detailed
        }
        
        # Get closest depth configuration
        depths = list(depth_configs.keys())
        closest_depth = min(depths, key=lambda x: abs(x - float(summary_depth)))
        config = depth_configs[closest_depth]
        
        return config

    def process_file(self, file_content, file_type, summary_depth=2.0):
        """Process file and return summary with format suggestions"""
        try:
            # Let HuggingFace handle the text extraction by treating file as raw input
            text = self._extract_text(file_content, file_type)
            
            # Get configuration based on summary depth
            config = self._optimize_length_params(len(text.split()), summary_depth)

            # Generate summary with error catching
            try:
                summary = self.summarizer(
                    text,
                    max_length=int(len(text.split()) * config['max_ratio']),
                    min_length=int(len(text.split()) * config['min_ratio']),
                    do_sample=False
                )[0]['summary_text']
            except Exception as e:
                logging.error(f"Summarization error: {str(e)}")
                summary = text[:int(len(text) * config['max_ratio'])]  # Fallback to truncation

            # Generate title suggestion based on summary
            suggested_title = self._generate_title(summary, config['title_length'])

            # Analyze content to suggest display format
            display_format = self._suggest_display_format(summary)

            return {
                'summary': summary,
                'title': suggested_title,
                'display_format': display_format
            }

        except Exception as e:
            logging.error(f"File processing error: {str(e)}")
            raise

    def _extract_text(self, file_content, file_type):
        """Extract text while preserving original file structure"""
        try:
            if file_type in ['pdf']:
                return extract_pdf_text(io.BytesIO(file_content))
            elif file_type in ['docx', 'doc']:
                doc = Document(io.BytesIO(file_content))
                return '\n'.join([p.text for p in doc.paragraphs])
            elif file_type in ['xlsx', 'xls']: 
                wb = load_workbook(io.BytesIO(file_content))
                text = []
                for sheet in wb.active:
                    for row in sheet.iter_rows(values_only=True):
                        text.append(' | '.join([str(cell) for cell in row if cell]))
                return '\n'.join(text)
            elif file_type in ['pptx', 'ppt']:
                prs = Presentation(io.BytesIO(file_content))
                text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                return '\n'.join(text)
            elif file_type in ['png', 'jpg', 'jpeg']:
                image = Image.open(io.BytesIO(file_content))
                return pytesseract.image_to_string(image)
            else:
                return file_content.decode('utf-8')
        except Exception as e:
            logging.error(f"Text extraction error: {str(e)}")
            return ""

    def _generate_title(self, text, max_words=4):
        """Generate a title from the summary text"""
        sentences = nltk.sent_tokenize(text)
        if not sentences:
            return "Document Summary"
            
        first_sentence = sentences[0]
        words = first_sentence.split()
        
        if len(words) <= max_words:
            title = ' '.join(words)
        else:
            title = ' '.join(words[:max_words])
            
        return title.strip().title()

    def _suggest_display_format(self, text):
        """Analyze text and suggest display format with visual elements"""
        sentences = nltk.sent_tokenize(text)
        words_per_sentence = len(text.split()) / len(sentences) if sentences else 0
        
        if words_per_sentence < 10:
            return {
                'type': 'bullet_points',
                'points': sentences,
                'style': {
                    'font': 'ComingSoon',
                    'colors': {
                        'primary': '#6A11CB',
                        'secondary': '#BC4E9C'
                    }
                },
                'image_query': 'concise notes visualization'
            }
        elif any(marker in text.lower() for marker in ['first', 'second', 'finally', 'next']):
            sections = self._extract_sections(text)
            return {
                'type': 'sections',
                'sections': sections,
                'style': {
                    'font': 'ComingSoon',
                    'colors': {
                        'primary': '#6A11CB',
                        'headers': '#BC4E9C'
                    }
                },
                'image_query': 'structured document organization'
            }
        else:
            return {
                'type': 'paragraph',
                'style': {
                    'font': 'ComingSoon',
                    'colors': {
                        'text': '#000000',
                        'background': '#FFFFFF'
                    }
                },
                'image_query': 'clean document layout'
            }

    def _extract_sections(self, text):
        """Extract sections from text with improved logic"""
        sentences = nltk.sent_tokenize(text)
        sections = []
        current_section = {'title': 'Overview', 'content': []}
        
        for sentence in sentences:
            lower_sentence = sentence.lower()
            is_new_section = any(marker in lower_sentence for marker in 
                               ['first', 'second', 'third', 'finally', 'next', 'moreover', 'furthermore'])
            
            if is_new_section:
                if current_section['content']:
                    sections.append({
                        'title': current_section['title'],
                        'content': ' '.join(current_section['content'])
                    })
                current_section = {
                    'title': sentence.strip(),
                    'content': []
                }
            else:
                current_section['content'].append(sentence)
                
        if current_section['content']:
            sections.append({
                'title': current_section['title'], 
                'content': ' '.join(current_section['content'])
            })
            
        return sections