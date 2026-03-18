import io
from PIL import Image
import pytesseract
from pdfminer.high_level import extract_text as extract_pdf_text
import docx
from pptx import Presentation
from openpyxl import load_workbook
import logging

class TextExtractor:
    @staticmethod
    def extract(file_content: bytes, file_type: str) -> str:
        file_type = file_type.lower()
        
        try:
            if file_type == 'pdf':
                return TextExtractor._extract_pdf(file_content)
            elif file_type in ['docx', 'doc']: # Note: doc may not work with python-docx, but we try
                return TextExtractor._extract_docx(file_content)
            elif file_type in ['pptx', 'ppt']:
                return TextExtractor._extract_pptx(file_content)
            elif file_type in ['xlsx', 'xls', 'csv']:
                return TextExtractor._extract_xlsx(file_content)
            elif file_type in ['png', 'jpg', 'jpeg', 'tiff', 'gif']:
                return TextExtractor._extract_image(file_content)
            else:
                return TextExtractor._extract_text(file_content)
        except Exception as e:
            logging.error(f"Failed to extract text from {file_type}: {str(e)}")
            # Fallback to simple decoding if specific parser fails
            return TextExtractor._extract_text(file_content)

    @staticmethod
    def _extract_pdf(file_content: bytes) -> str:
        pdf_file = io.BytesIO(file_content)
        return extract_pdf_text(pdf_file)

    @staticmethod
    def _extract_docx(file_content: bytes) -> str:
        doc_file = io.BytesIO(file_content)
        doc = docx.Document(doc_file)
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])

    @staticmethod
    def _extract_pptx(file_content: bytes) -> str:
        ppt_file = io.BytesIO(file_content)
        prs = Presentation(ppt_file)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    @staticmethod
    def _extract_xlsx(file_content: bytes) -> str:
        excel_file = io.BytesIO(file_content)
        wb = load_workbook(excel_file, data_only=True)
        text = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join([str(cell) for cell in row if cell is not None])
                if row_text:
                    text.append(row_text)
        return "\n".join(text)

    @staticmethod
    def _extract_image(file_content: bytes) -> str:
        image = Image.open(io.BytesIO(file_content))
        return pytesseract.image_to_string(image)

    @staticmethod
    def _extract_text(file_content: bytes) -> str:
        try:
            return file_content.decode('utf-8')
        except UnicodeDecodeError:
            try:
                return file_content.decode('latin-1')
            except Exception:
                return ""
